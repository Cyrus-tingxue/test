import os
import io
import re
import json
import time
import uuid
import tempfile
import traceback
import subprocess
import sys
from fastapi import APIRouter, HTTPException, UploadFile, File, Form
from fastapi.responses import StreamingResponse, FileResponse, JSONResponse

from schemas import CreativeRequest, CodeRequest, PPTRequest, PPTCreateRequest, MindMapRequest
from core_helpers import resolve_api_key, resolve_model, keepalive_llm_stream, TEMP_DOWNLOADS
from utils.llm_client import acall_llm_stream
from utils.file_parser import parse_file

router = APIRouter()

@router.post("/creative")
async def generate_creative(request: CreativeRequest):
    try:
        prompt = ""
        f = request.fields
        extra = {}
        if request.base_url: extra["api_base"] = request.base_url
        
        if request.task == "daily_report":
            prompt = f"请根据以下工作内容生成一份专业的日报：\n{f.get('content', '')}\n\n要求：条理清晰，分为【今日工作】、【存在问题】、【明日计划】三部分。"
        elif request.task == "weekly_report":
            prompt = f"请根据以下工作内容生成一份周报：\n{f.get('content', '')}\n\n要求：总结本周重点，分析数据/成果，规划下周。"
        elif request.task == "email":
            prompt = f"请帮我写一封邮件。\n收件人：{f.get('receiver', '')}\n主题：{f.get('topic', '')}\n要点：{f.get('content', '')}\n\n要求：语气商务、得体。"
        elif request.task == "translation":
            prompt = f"请将以下内容翻译成 {f.get('target_lang', '中文')}，并进行润色：\n{f.get('content', '')}"
        elif request.task == "meeting_minutes":
            prompt = f"请根据以下会议纪要草稿，整理成一份正式的会议纪要：\n{f.get('content', '')}\n\n要求：包含会议主题、时间、参会人员、主要决议、待办事项。"
        elif request.task == "okr_draft":
            prompt = f"请根据以下目标，帮我起草一份OKR（目标与关键结果）：\n目标：{f.get('content', '')}\n\n要求：符合SMART原则，包含1个O和3-5个KR。"
        elif request.task == "essay_outline":
            prompt = f"请为以下论文/文章主题生成一份详细大纲：\n主题：{f.get('content', '')}\n\n要求：结构严谨，包含引言、各章节论点、结论。"
        elif request.task == "study_plan":
            prompt = f"请为我制定一份学习计划。\n学习科目/技能：{f.get('topic', '')}\n可用时间：{f.get('time', '')}\n\n要求：分阶段，可执行性强。"
        elif request.task == "xhs_copy":
            prompt = f"请写一篇小红书风格的种草文案。\n产品/主题：{f.get('topic', '')}\n卖点/内容：{f.get('content', '')}\n\n要求：标题吸引人（带emoji），正文亲切活泼，包含Tag。"
        elif request.task == "video_script":
            prompt = f"请写一份短视频脚本。\n主题：{f.get('topic', '')}\n\n要求：包含分镜描述、台词、画面建议，时长约1分钟。"
        elif request.task == "recipe_gen":
            prompt = f"请根据以下食材生成一份菜谱：\n食材：{f.get('content', '')}\n\n要求：包含菜名、所需配料、详细烹饪步骤。"
        elif request.task == "travel_plan":
            prompt = f"请为我制定一份旅行计划。\n目的地：{f.get('destination', '')}\n天数：{f.get('days', '')}\n\n要求：包含每日行程安排、景点推荐、美食建议。"
        elif request.task == "resume_polish":
            prompt = f"请帮我优化这份简历内容：\n{f.get('content', '')}\n\n要求：使用专业职场术语，突出成就和数据，优化排版建议。"
        elif request.task == "interview_prep":
            prompt = f"即将面试岗位：{f.get('topic', '')}\n\n请列出5个高频面试题，并给出优秀的回答思路（STAR法则）。"
        elif request.task == "swot_analysis":
            prompt = f"请对以下项目/主题进行SWOT分析：\n{f.get('content', '')}\n\n要求：列出优势(S)、劣势(W)、机会(O)、威胁(T)，并给出战略建议。"
        elif request.task == "contract_review":
            prompt = f"我是乙方/个人，请帮我审查虽然条款，指出潜在风险和陷阱：\n{f.get('content', '')}\n\n要求：通俗易懂，标出高风险条款。"
        elif request.task == "title_gen":
            prompt = f"请为这篇文章生成10个爆款标题：\n内容/主题：{f.get('content', '')}\n目标受众：{f.get('topic', '')}\n\n要求：包含数字、悬念、痛点等爆款元素。"
        elif request.task == "article_polish":
            prompt = f"请润色以下文章，使其更流畅、更有文采：\n{f.get('content', '')}\n\n要求：修正语病，提升词汇丰富度，保持原意。"
        elif request.task == "excel_gen":
            prompt = f"请根据以下描述生成一份CSV格式的数据，用于Excel处理：\n描述：{f.get('content', '')}\n\n要求：只输出CSV内容，不要有任何解释，首行为表头。"
            try:
                import pandas as pd
                def process_excel(csv_content):
                    if "```csv" in csv_content:
                        csv_content = csv_content.split("```csv")[1].split("```")[0].strip()
                    elif "```" in csv_content:
                        csv_content = csv_content.split("```")[1].split("```")[0].strip()
                    df = pd.read_csv(io.StringIO(csv_content), on_bad_lines='warn', sep=None, engine='python')
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                        tmp_path = tmp.name
                    df.to_excel(tmp_path, index=False)
                    file_id = str(uuid.uuid4())
                    TEMP_DOWNLOADS[file_id] = {
                        "path": tmp_path,
                        "filename": f"generated_excel_{int(time.time())}.xlsx",
                        "media_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        "ts": time.time()
                    }
                    return {"download_url": f"/api/download/{file_id}"}
                
                return StreamingResponse(
                    keepalive_llm_stream(
                        provider=request.provider,
                        model=resolve_model(request.model),
                        api_key=resolve_api_key(request.api_key),
                        messages=[{"role": "user", "content": prompt}],
                        process_fn=process_excel,
                        **extra
                    ),
                    media_type="application/json"
                )
            except Exception as e:
                import traceback
                traceback.print_exc()
                raise HTTPException(status_code=500, detail=f"Excel Gen Failed: {str(e)}")
        else:
            prompt = f"请完成以下任务：{f.get('content', '')}"

        async def event_generator():
            try:
                system_instruction = "你是一个专业的AI内容生成助手。请直接根据用户的要求生成最终的内容草稿，不要追问细节，不要反问。如果信息不足，请自行补充合理的假设内容以完成生成。"
                messages = [
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": prompt}
                ]
                async for chunk in acall_llm_stream(
                    provider=request.provider,
                    model=resolve_model(request.model),
                    api_key=resolve_api_key(request.api_key),
                    messages=messages,
                    **extra
                ):
                    yield chunk
            except Exception as e:
                yield f"Error: {str(e)}"

        return StreamingResponse(event_generator(), media_type="text/plain")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/creative/upload")
async def generate_creative_upload(
    file: UploadFile = File(...),
    task: str = Form(...),
    content: str = Form(""),
    provider: str = Form("OpenRouter"),
    model: str = Form("gpt-3.5-turbo"),
    api_key: str = Form(""),
    base_url: str = Form(None),
    receiver: str = Form(None),
    topic: str = Form(None),
    target_lang: str = Form(None),
    time: str = Form(None),
    destination: str = Form(None),
    days: str = Form(None)
):
    try:
        file_content = ""
        try:
            file_bytes = await file.read()
            class NamedBytesIO(io.BytesIO):
                def __init__(self, content, name):
                    super().__init__(content)
                    self.name = name
            f_obj = NamedBytesIO(file_bytes, file.filename)
            file_content = parse_file(f_obj)
        except Exception as e:
            return JSONResponse(status_code=400, content={"detail": f"File Parse Error: {str(e)}"})

        full_content = f"{content}\n\n【附件内容 ({file.filename})】:\n{file_content}"
        fields = {
            "content": full_content,
            "receiver": receiver,
            "topic": topic,
            "target_lang": target_lang,
            "time": time,
            "destination": destination,
            "days": days
        }
        fields = {k: v for k, v in fields.items() if v is not None}
        req = CreativeRequest(
            task=task,
            fields=fields,
            provider=provider,
            model=model,
            api_key=resolve_api_key(api_key),
            base_url=base_url
        )
        return await generate_creative(req)
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/ppt/outline")
async def generate_ppt_outline(request: PPTRequest):
    try:
        _PPT_PROMPT = """你是一个专业的 PPT 设计师。请根据用户的主题："{topic}" 生成一份 PPT 大纲，包含 5-8 页内容。
        请严格返回合法的 JSON 格式，结构如下：
        {{"slides": [{{"title": "页标题", "content": ["要点1"]}}]}}
        1. 第一页必须是封面页（"title": "主题", "content": ["副标题/演讲者"]）。
        2. 内容要精炼、专业。不要包含 JSON 以外的任何文字。"""
        prompt = _PPT_PROMPT.format(topic=request.topic)
        messages = [{"role": "user", "content": prompt}]
        extra = {}
        if request.base_url: extra["api_base"] = request.base_url

        def process_ppt(response_text):
            clean_json = re.sub(r"```json\s*", "", response_text)
            clean_json = re.sub(r"\s*```", "", clean_json)
            data = json.loads(clean_json)
            return {"data": data}

        return StreamingResponse(
            keepalive_llm_stream(
                provider=request.provider,
                model=resolve_model(request.model),
                api_key=resolve_api_key(request.api_key),
                messages=messages,
                process_fn=process_ppt,
                **extra
            ),
            media_type="application/json"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/ppt/create")
async def create_ppt(request: PPTCreateRequest):
    try:
        from pptx import Presentation
        data = request.data
        slides_data = data.get("slides", [])
        prs = Presentation()
        if slides_data:
            layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            cover = slides_data[0]
            title.text = cover.get("title", "Presentation")
            subtitle.text = "\n".join(cover.get("content", []))
        layout = prs.slide_layouts[1]
        for slide_info in slides_data[1:]:
            slide = prs.slides.add_slide(layout)
            shapes = slide.shapes
            title_shape = shapes.title
            title_shape.text = slide_info.get("title", "Untitled")
            body_shape = shapes.placeholders[1]
            tf = body_shape.text_frame
            points = slide_info.get("content", [])
            if points:
                tf.text = points[0]
                for point in points[1:]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 1
        ppt_buffer = io.BytesIO()
        prs.save(ppt_buffer)
        ppt_buffer.seek(0)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt:
            tmp_ppt.write(ppt_buffer.getvalue())
            tmp_path = tmp_ppt.name
        return FileResponse(tmp_path, media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation", filename="generated_presentation.pptx")
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/mindmap")
async def generate_mindmap(request: MindMapRequest):
    try:
        prompts = {
            "mindmap": """你是一个思维导图专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的思维导图代码 (mindmap)。
                【输出规则】必须以 `mindmap` 关键字开头。根节点格式：`root((主题名))`。使用空格缩进表示层级。仅返回代码本身，禁止包含任何 Markdown 标记。至少扩展 3 层，包含 10-15 个节点。节点文字包含特殊字符时请用双引号包裹。""",
            "flowchart": """你是一个流程图专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的流程图代码 (graph TD)。
                【输出规则】必须包含分支判断逻辑(Decision Nodes)使用`{{}}`语法。必须是非线性的。节点文字若包含特殊内容需用双引号包裹`id["文字"]`。至少10-15个不同节点。仅返回Mermaid代码。""",
            "timeline": """你是一个时间序列专家。请根据用户的主题："{topic}"生成 Mermaids JS timeline。包含 5-8 个关键点。仅放代码。""",
            "gantt": """你设计甘特图。主题："{topic}" 仅输出包含 dateFormat 的 gantt 代码。""",
            "sequence": """你设计时序图 (sequenceDiagram)。主题："{topic}" 至少 3 参与者交互代码。""",
            "class": """你写架构类图 (classDiagram)。主题："{topic}" 仅输出代码。""",
            "state": """你设计状态转换 (stateDiagram-v2)。主题："{topic}" 仅输出代码。""",
            "pie": """你设计饼图 (pie)。主题："{topic}" 标签必须用双引号包裹，例如："分类A" : 40。仅输出代码。"""
        }
        selected_prompt = prompts.get(request.chart_type, prompts["mindmap"])
        prompt = selected_prompt.format(topic=request.topic)
        messages = [{"role": "user", "content": prompt}]
        extra = {}
        if request.base_url: extra["api_base"] = request.base_url

        def process_mindmap(response_text):
            code = response_text.strip()
            blocks = re.findall(r"```(?:mermaid)?\s*(.*?)\s*```", code, re.DOTALL)
            if blocks: code = blocks[0].strip()
            start_keywords = ["mindmap", "graph", "flowchart", "timeline", "gantt", "sequenceDiagram", "classDiagram", "stateDiagram", "pie"]
            lines = code.split('\n')
            valid_start = -1
            for i, line in enumerate(lines):
                clean_line = line.strip()
                if any(clean_line.startswith(kw) for kw in start_keywords):
                    valid_start = i; break
            if valid_start != -1: code = "\n".join(lines[valid_start:]).strip()
            if request.chart_type == "pie":
                if not code.startswith("pie"): code = "pie\n" + code
                code = re.sub(r'^\s*([^"\s:][^:\n]+[^"\s:])\s*:\s*(\d+)', r'    "\1" : \2', code, flags=re.MULTILINE)
            if request.chart_type == "mindmap" and not code.startswith("mindmap"):
                code = "mindmap\n  " + code.replace("\n", "\n  ")
            return {"code": code}

        return StreamingResponse(keepalive_llm_stream(provider=request.provider, model=resolve_model(request.model), api_key=resolve_api_key(request.api_key), messages=messages, process_fn=process_mindmap, **extra), media_type="application/json")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
