"""Office AI Mate - Backend API Server"""

from fastapi.responses import StreamingResponse, FileResponse, JSONResponse
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Request
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
import uvicorn
import os
import json
import tempfile
import io
import re
import uuid
import time
import random
import asyncio
import pandas as pd
from typing import Any
import subprocess
import sys
import traceback
try:
    from curl_cffi.requests import AsyncSession
    curl_cffi_installed = True
except ImportError:
    curl_cffi_installed = False   
    AsyncSession = None   
from utils.llm_client import call_llm, call_llm_stream, acall_llm, acall_llm_stream
from utils.file_parser import parse_file
from utils.chat_analyzer import analyze_chat_style
from utils.local_scanner import scan_qq_logs

# ── Default config from environment variables (keep API key on server side) ──
DEFAULT_PROVIDER = os.environ.get("DEFAULT_PROVIDER", "OpenRouter")
DEFAULT_MODEL = os.environ.get("DEFAULT_MODEL", "arcee-ai/trinity-large-preview:free")
DEFAULT_API_KEY = os.environ.get("DEFAULT_API_KEY", "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")

def resolve_api_key(key: str) -> str:
    """当前端未提供 API Key 时，回退到默认值。"""
    return key.strip() if key and key.strip() else DEFAULT_API_KEY

def resolve_model(model: str) -> str:
    """当前端未提供模型名时，回退到默认模型。"""
    return model.strip() if model and model.strip() else DEFAULT_MODEL

# ── 连接保活辅助 (解决中国用户访问美国服务器时非流式 API 超时断连) ──

# 临时文件下载缓存: {file_id: {"path": str, "filename": str, "media_type": str, "ts": float}}
TEMP_DOWNLOADS = {}

async def keepalive_llm_stream(provider, model, api_key, messages, process_fn=None, **extra):
    """空格保活的流式 LLM 调用生成器。

    在 LLM 处理期间每 3 秒发送一个空格字符保持 TCP 连接活跃。
    LLM 完成后，调用 process_fn 处理完整文本，最终发送 JSON 结果。
    JSON.parse 会自动忽略前导空格，所以前端 safeJson() 无需任何改动。
    """
    collected = []
    done = asyncio.Event()
    error_holder = {}

    async def llm_task():
        try:
            async for chunk in acall_llm_stream(
                provider=provider,
                model=model,
                api_key=api_key,
                messages=messages,
                **extra
            ):
                collected.append(chunk)
        except Exception as e:
            error_holder['error'] = e
        finally:
            done.set()

    task = asyncio.create_task(llm_task())

    # 每 3 秒发送空格保活
    while not done.is_set():
        yield " "
        try:
            await asyncio.wait_for(asyncio.shield(done.wait()), timeout=3.0)
        except asyncio.TimeoutError:
            pass

    # LLM 完成，发送最终结果
    if 'error' in error_holder:
        yield json.dumps({"detail": str(error_holder['error'])})
    else:
        full_text = "".join(collected)
        if process_fn:
            try:
                result = process_fn(full_text)
                yield json.dumps(result, ensure_ascii=False)
            except Exception as e:
                yield json.dumps({"detail": f"处理结果时出错: {str(e)}"})
        else:
            yield json.dumps({"text": full_text}, ensure_ascii=False)

# Initialize FastAPI app
app = FastAPI(title="Office AI Mate API", version="2.0", description="By 昨夜提灯看雪")

# Enable CORS (Cross-Origin Resource Sharing)
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Allow all origins for development
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global Exception Handler
@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    print(f"Global Exception: {exc}")
    return JSONResponse(
        status_code=500,
        content={"detail": f"Internal Server Error: {str(exc)}"},
    )

# Request Models
class SearchRequest(BaseModel):
    query: str
    page: int = 1
    max_results: int = 8  # Reduced to 8 as per user requestol = True
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo" 
    api_key: str = ""
    base_url: str | None = None

class ChatRequest(BaseModel):
    messages: list[dict]
    provider: str = "OpenRouter"
    model: str = "google/gemini-2.0-flash-exp:free"
    api_key: str = ""
    base_url: str | None = None

class PPTRequest(BaseModel):
    topic: str
    items: dict = {}
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None

class PPTCreateRequest(BaseModel):
    data: dict

class CreativeRequest(BaseModel):
    task: str
    fields: dict
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None

class CodeRequest(BaseModel):
    task: str
    language: str
    content: str
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None



class SystemRequest(BaseModel):
    query: str
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None

class GameRequest(BaseModel):
    pet_state: dict
    action: str
    user_input: str = ""
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None

class SystemExecuteRequest(BaseModel):
    code: str


# --- Global Search Cache ---
# Format: { "query_str": { "ts": timestamp, "results": [list_of_dicts] } }
SEARCH_CACHE = {}
CACHE_TTL = 600 # 10 minutes

# --- API Endpoints ---

@app.get("/api/health")
async def health_check():
    return {"status": "ok", "message": "Office AI Mate API is running"}

@app.get("/api/download/{file_id}")
async def download_temp_file(file_id: str):
    """下载由保活端点生成的临时文件。"""
    # 清理过期文件 (超过 10 分钟)
    expired = [k for k, v in TEMP_DOWNLOADS.items() if time.time() - v['ts'] > 600]
    for k in expired:
        path = TEMP_DOWNLOADS.pop(k, {}).get('path', '')
        if path and os.path.exists(path):
            try: os.remove(path)
            except: pass

    info = TEMP_DOWNLOADS.pop(file_id, None)
    if not info or not os.path.exists(info['path']):
        raise HTTPException(status_code=404, detail="File not found or expired")
    return FileResponse(info['path'], media_type=info['media_type'], filename=info['filename'])

@app.get("/api/defaults")
async def get_defaults():
    """返回默认配置（不暴露 API Key）。"""
    return {
        "provider": DEFAULT_PROVIDER,
        "model": DEFAULT_MODEL,
        "has_default_key": bool(DEFAULT_API_KEY)
    }

@app.post("/api/search")
async def search(request: SearchRequest, raw_request: Request):
    print(f"INFO: Search Request: {request.query} (Page {request.page})")
    print(f"DEBUG: Processing Page {request.page} with Max Results {request.max_results}")
    optimized_query = request.query
    
    # --- Cache Lookup ---
    # Clean old cache
    current_time = time.time()
    # Create list of keys to remove first to avoid runtime error
    keys_to_remove = [k for k, v in SEARCH_CACHE.items() if current_time - v['ts'] > CACHE_TTL]
    for k in keys_to_remove:
        del SEARCH_CACHE[k]

    cached_data = SEARCH_CACHE.get(optimized_query)
    
    # If loading more (page > 1) and we have cache, try to serve from it
    if request.page > 1 and cached_data:
        start_idx = (request.page - 1) * request.max_results
        end_idx = start_idx + request.max_results
        
        # If we have enough items in cache, return them
        if len(cached_data['results']) >= end_idx:
            print(f"INFO: Serving Page {request.page} from CACHE ({len(cached_data['results'])} items cached)")
            return {"results": cached_data['results'][start_idx:end_idx], "optimized_query": optimized_query}
        else:
             print(f"INFO: Cache exhausted for Page {request.page} (only {len(cached_data['results'])} items). Will try to fetch more.")
             # We let it fall through to fetch more, but we should carry over existing results to filter duplicates?
             # Simple approach: just try to fetch new items and APPEPEND to cache?
             pass 

    results = []
    errors = []
    seen_urls = set()
    
    # Pre-fill seen_urls from cache if expanding
    if cached_data:
        for r in cached_data['results']:
            seen_urls.add(r['href'])

    # helper
    def add_result(title, href, body):
        if not href or href in seen_urls: return
        # Filter relative links that accidentally got through
        if not href.startswith("http"): return 
        
        # Deduplicate by title similarity (simple)
        for r in results:
            if r['title'] == title.strip(): return
        
        # Also check against cache if we are extending it
        if cached_data:
             for r in cached_data['results']:
                 if r['title'] == title.strip(): return

        seen_urls.add(href)
        item = {
            "title": title.strip(),
            "href": href.strip(),
            "body": body.strip()
        }
        results.append(item)
        if cached_data:
             cached_data['results'].append(item)

    # --- 0. Dependency Check ---
    global curl_cffi_installed
    global AsyncSession
    httpx_module = None
    bs4_module = None
    ddgs_module = None 
    
    # Try Import All
    try:
        import httpx
        from bs4 import BeautifulSoup
        from duckduckgo_search import DDGS
        httpx_module = httpx
        bs4_module = BeautifulSoup
        ddgs_module = DDGS
        
        # If curl_cffi was already installed but not yet imported in this thread
        if not curl_cffi_installed:
            try:
                import curl_cffi.requests
                AsyncSession = curl_cffi.requests.AsyncSession
                curl_cffi_installed = True
            except ImportError:
                pass
                
    except ImportError:
        pass

    # Forced Install if anyone is missing
    if not (httpx_module and bs4_module and curl_cffi_installed and ddgs_module):
        print("DEBUG: Missing critical dependencies, attempting auto-install...")
        try:
            import sys
            import subprocess
            print("INFO: Installing curl_cffi, httpx, beautifulsoup4, duckduckgo-search...")
            subprocess.check_call([sys.executable, "-m", "pip", "install", "curl_cffi", "httpx", "beautifulsoup4", "duckduckgo-search"])
            import httpx
            from bs4 import BeautifulSoup
            import curl_cffi.requests
            from duckduckgo_search import DDGS
            
            httpx_module = httpx
            bs4_module = BeautifulSoup
            AsyncSession = curl_cffi.requests.AsyncSession
            curl_cffi_installed = True
            ddgs_module = DDGS
        except Exception as install_e:
            err_msg = f"Dependency Install Failed: {install_e}"
            print(f"ERROR: {err_msg}")
            errors.append(err_msg)


    # Random User-Agent Pool (High Quality Real UAs)
    USER_AGENTS_POOL = [
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/121.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36",
        "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/605.1.15 (KHTML, like Gecko) Version/17.2 Safari/605.1.15",
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64; rv:123.0) Gecko/20100101 Firefox/123.0",
        "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/122.0.0.0 Safari/537.36"
    ]

    # Base Headers
    BASE_HEADERS = {
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,image/apng,*/*;q=0.8",
        "Accept-Language": "zh-CN,zh;q=0.9,en;q=0.8",
        "Accept-Encoding": "gzip, deflate, br",
        "DNT": "1",
        "Connection": "keep-alive",
        "Upgrade-Insecure-Requests": "1",
        "User-Agent": USER_AGENTS_POOL[0]  # Ensure a default UA is present
    }
    
    # --- Helper: Generic Parser ---
    def generic_parse(html, source_name):
        soup = bs4_module(html, "html.parser")
        found = []
        
        # 1. Try Common Search Result Containers
        # Baidu: .result, .c-container, .new-pmd
        # Sogou: .vrwrap, .rb, .result
        # 360: .res-list, .result
        # Bing: .b_algo
        potential_items = soup.select(".result, .c-result, .c-container, .vrwrap, .rb, .res-list, .b_algo, .lib-box, li")
        
        for item in potential_items:
            # Title is usually in h3 > a or h2 > a
            title_tag = item.select_one("h3 a, h2 a, h3, h2")
            if not title_tag: continue
            
            # If title tag is h3/h2 itself, check if it has 'a' child, or text
            if title_tag.name in ['h3', 'h2']:
                if title_tag.find('a'):
                    title_tag = title_tag.find('a')
                else:
                    # Sometimes the h3 itself is usable? No, usually needs a link.
                    continue

            link = title_tag.get("href", "")
            if not link or link.startswith("javascript") or link == "#": continue
            
            title = title_tag.get_text(strip=True)
            if not title: continue

            # Body/Snippet
            snippet_cat = item.select_one(".c-abstract, .text-layout, .res-desc, .b_caption p, .b_snippet, .content-right_8Zs40, .ft, p")
            body = snippet_cat.get_text(strip=True) if snippet_cat else ""
            
            # Fix Relative Links
            if source_name == "Sogou" and link.startswith("/link"): link = "https://www.sogou.com" + link
            if source_name == "Baidu" and link.startswith("/"): link = "https://www.baidu.com" + link
            
            found.append({'title': title, 'href': link, 'body': body})
            
        return found

    # --- 1. Baidu (Priority) ---
    if bs4_module:
        try:
            time.sleep(random.uniform(0.5, 1.5))
            page = request.page
            params: dict[str, Any] = {"wd": optimized_query}
            headers = BASE_HEADERS.copy()
            headers["Host"] = "www.baidu.com"
            if page > 1: params["pn"] = (page - 1) * 10
            
            resp_text = ""
            if curl_cffi_installed:
                try:
                    async with AsyncSession(impersonate="chrome120", verify=False, timeout=10) as client:
                        resp = await client.get("https://www.baidu.com/s", params=params, headers=headers)
                        resp_text = resp.text
                except Exception as cf_e:
                    errors.append(f"Baidu (curl_cffi error): {str(cf_e)}")
            
            # Fallback to httpx if curl_cffi failed or not installed
            if not resp_text and httpx_module:
                headers["User-Agent"] = random.choice(USER_AGENTS_POOL)
                async with httpx_module.AsyncClient(timeout=6, verify=False) as client:
                    resp = await client.get("https://www.baidu.com/s", params=params, headers=headers)
                    resp_text = resp.text

            if len(resp_text) > 2000:
                items = generic_parse(resp_text, "Baidu")
                if items:
                    for item in items:
                        add_result(item['title'], item['href'], item['body'])
                        if len(results) >= request.max_results: break
                else:
                    errors.append(f"Baidu: Parsed 0 items (Len: {len(resp_text)})")
            elif resp_text:
                snippet = resp_text[:100].replace("\n", " ")
                errors.append(f"Baidu: Blocked (Len: {len(resp_text)}, Content: {snippet}...)")
            else:
                errors.append("Baidu: No response text")

        except Exception as e:
            errors.append(f"Baidu (General): {str(e)}")




    # --- 2. Bing Global (Fallback) ---
    if len(results) < request.max_results and bs4_module:
        try:
            page = int(request.page)
            first = (page - 1) * 10 + 1
            params = {"q": optimized_query, "first": first, "count": 10, "setmkt": "en-US"}
            print(f"DEBUG: Bing Params: {params}")
            headers = BASE_HEADERS.copy()
            # headers["User-Agent"] = ... # Already in BASE_HEADERS
            
            resp_text = ""
            if curl_cffi_installed:
                try:
                    async with AsyncSession(impersonate="chrome120", verify=False, timeout=15) as client:
                        resp = await client.get("https://www.bing.com/search", params=params, headers=headers)
                        resp_text = resp.text
                except Exception as cf_e:
                     errors.append(f"Bing (curl_cffi error): {str(cf_e)}")

            if not resp_text and httpx_module:
                headers["User-Agent"] = random.choice(USER_AGENTS_POOL)
                async with httpx_module.AsyncClient(timeout=10, verify=False) as client:
                    resp = await client.get("https://www.bing.com/search", params=params, headers=headers)
                    resp_text = resp.text

            if resp_text:
                items = generic_parse(resp_text, "Bing")
                for item in items:
                    add_result(item['title'], item['href'], item['body'])
                    if len(results) >= request.max_results: break
                if not results:
                     errors.append(f"Bing: No results (Len: {len(resp_text)})")
            else:
                errors.append("Bing: No response text")
                 
        except Exception as e:
            errors.append(f"Bing: {str(e)}")

    # --- 3. DuckDuckGo (Fallback) ---
    # DDGS has multiple backends (api, html, lite). If one fails (e.g. rate limit), try others.
    # We fetch a large batch (60 items) to cover pages 1-6 if possible.
    
    start_index = (request.page - 1) * request.max_results
    # If page 1, we try to fetch 60 (cache warm up). If page > 1, we fetch what we need.
    needed_count = 60 if request.page == 1 else (start_index + request.max_results + 5)
    
    if ddgs_module and len(results) < request.max_results:
        ddg_backends = ['api', 'html', 'lite']
        ddg_success = False
        
        for backend in ddg_backends:
            if ddg_success: break
            try:
                print(f"INFO: Trying DuckDuckGo (Backend: {backend}) for '{optimized_query}' (Need {needed_count})...")
                with ddgs_module() as ddgs:
                    ddgs_results = ddgs.text(
                        optimized_query, 
                        region='wt-wt', 
                        safesearch='moderate', 
                        timelimit='y', 
                        max_results=needed_count,
                        backend=backend
                    )
                    
                    if ddgs_results:
                        valid_items_found = 0
                        # Manually slice for pagination correctness
                        current_idx = 0
                        
                        # Store all results in cache if page 1
                        all_found_items = []

                        for r in ddgs_results:
                            item_title = r.get('title', '')
                            item_href = r.get('href', '')
                            item_body = r.get('body', '')
                            
                            if not item_href: continue

                            # Collect everything for cache
                            all_found_items.append({
                                "title": item_title,
                                "href": item_href,
                                "body": item_body
                            })

                            # Process for current page return
                            current_idx += 1
                            if current_idx <= start_index: continue # Skip previous pages
                            
                            add_result(item_title, item_href, item_body)
                            if len(results) >= request.max_results: 
                                # Don't break here, we want to finish collecting for cache if page 1? 
                                # Actually, add_result adds to 'results' list which is returned.
                                # But we also have 'SEARCH_CACHE' logic below.
                                # To keep it simple: just fill 'results' for NOW, and let existing Cache Logic handle the rest.
                                pass 
                        
                        if results:
                            ddg_success = True
                            print(f"INFO: DuckDuckGo ({backend}) success. Found {len(results)} items for current page.")
                            
                        # If we fetched a lot (Page 1), ensure we pass them to the cache logic
                        # The existing cache logic uses `results` list. 
                        # IF we are on page 1, `add_result` will populate `results`.
                        # But `add_result` stops adding to `results` once `request.max_results` is reached due to the check I removed above?
                        # Wait, I removed the break above to allow collecting more for cache?
                        # No, `add_result` doesn't check limit. The calling loop usually checked `if len(results) >= max: break`.
                        
            except Exception as ddg_e:
                 print(f"WARN: DuckDuckGo ({backend}) failed: {ddg_e}")
                 errors.append(f"DDG-{backend}: {str(ddg_e)}")
    
    # Prune results to max_results if we over-fetched for cache purposes
    # (Only needed if I removed the break logic)
    # Actually, let's keep the break logic in the loop for the *response* results, 
    # but we might want to populate cache with MORE.
    # The current Architecture makes it hard to decouple "response" from "cache population" without big refactor.
    # Let's stick to: "Fetch what we need, return what we need". 
    # The cache will naturally build up if we fetch 60 and add them to `results`? 
    # NO, `results` is what is returned to user.
    
    # Correct Logic:
    # 1. We have `results` list which goes to user.
    # 2. We have `SEARCH_CACHE` which stores *everything*.
    # 3. `add_result` appends to BOTH `results` (unconditionally?) and `SEARCH_CACHE['results']`.
    # 4. Wait, `add_result` does: `results.append(...)`, `if cached_data: cached_data['results'].append(...)`.
    
    # So if we want to populate cache with 60 items but only return 10 to user:
    # We should NOT break the loop. We should iterate all 60.
    # BUT `results` will grow to 60.
    # So at the END of the function, we must slice `results` before returning.
    # The existing code at the bottom does: `final_results = results[:request.max_results]` if Page 1.
    # PERFECT. So we just need to NOT break the loop.


    # --- 5. Error Reporting ---
    if not results:
        all_errors = "; ".join(errors)
        print(f"ERROR: All search attempts failed. {all_errors}")
        
        friendly_body = "搜索服务暂不可用 (所有引擎均失败)"
        if "No results" in all_errors:
             friendly_body += " - 访问成功但无法解析内容 (反爬/验证码)"
        if "timeout" in all_errors or "ConnectError" in all_errors:
             friendly_body += " - 网络连接问题"
             
        results.append({
            "title": "❌ 搜索失败",
            "href": "#",
            "body": f"{friendly_body}\n调试日志: {all_errors}"
        })

    # --- Cache Save (New Search) ---
    if request.page == 1 and results:
        SEARCH_CACHE[optimized_query] = {
            "ts": time.time(),
            "results": results # these are all results fetched (up to 50)
        }
        # Return only the requested amount for the first page
        final_results = results[:request.max_results]
        return {"results": final_results, "original_query": request.query, "optimized_query": optimized_query}
    
    # --- Cache Update (Appended) ---
    if request.page > 1 and results:
        # results now contains ONLY new items found during this fetch (because add_result filters against cache)
        # Note: add_result ALREADY appended to cached_data['results'] in place if cached_data existed
        # So we just need to return the new items.
         return {"results": results[:request.max_results], "original_query": request.query, "optimized_query": optimized_query}

    return {"results": results, "original_query": request.query, "optimized_query": optimized_query}

@app.post("/api/chat")
async def chat(request: ChatRequest):
    try:
        api_key = resolve_api_key(request.api_key)
        if not api_key:
             raise HTTPException(status_code=400, detail="API Key is required")

        extra = {}
        if request.base_url:
            extra["api_base"] = request.base_url
            
        async def event_generator():
            try:
                async for chunk in acall_llm_stream(
                    provider=request.provider,
                    model=resolve_model(request.model),
                    api_key=api_key,
                    messages=request.messages,
                    **extra
                ):
                    yield chunk
            except Exception as e:
                yield f"Error: {str(e)}"

        return StreamingResponse(event_generator(), media_type="text/plain")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate/creative")
async def generate_creative(request: CreativeRequest):
    try:
        # Prompt Engineering based on Task
        prompt = ""
        f = request.fields
        
        
        extra = {}
        if request.base_url: extra["api_base"] = request.base_url
        
        if request.task == "daily_report":
            prompt = f"请根据以下工作内容生成一份专业的日报：\n{f.get('content', '')}\n\n要求：条理清晰，分为【今日工作】、【存在问题】、【明日计划】三部分。"
        elif request.task == "weekly_report":
            prompt = f"请根据以下工作内容生成一份周报：\n{f.get('content', '')}\n\n要求：总结本周重点，分析数据/成果，规划下周。"
        elif request.task == "email":
            prompt = f"请帮我写一封邮件。\n收件人：{f.get('receiver', '')}\n主题：{f.get('topic', '')}\n要点：{f.get('content', '')}\n\n要求：语气商务、得体。"
        elif request.task == "translation":
            prompt = f"请将以下内容翻译成 {f.get('target_lang', '中文')}，并进行润色：\n{f.get('content', '')}"
        
        # --- Work ---
        elif request.task == "meeting_minutes":
            prompt = f"请根据以下会议纪要草稿，整理成一份正式的会议纪要：\n{f.get('content', '')}\n\n要求：包含会议主题、时间、参会人员、主要决议、待办事项。"
        elif request.task == "okr_draft":
            prompt = f"请根据以下目标，帮我起草一份OKR（目标与关键结果）：\n目标：{f.get('content', '')}\n\n要求：符合SMART原则，包含1个O和3-5个KR。"
            
        # --- Education ---
        elif request.task == "essay_outline":
            prompt = f"请为以下论文/文章主题生成一份详细大纲：\n主题：{f.get('content', '')}\n\n要求：结构严谨，包含引言、各章节论点、结论。"
        elif request.task == "study_plan":
            prompt = f"请为我制定一份学习计划。\n学习科目/技能：{f.get('topic', '')}\n可用时间：{f.get('time', '')}\n\n要求：分阶段，可执行性强。"

        # --- Social Media ---
        elif request.task == "xhs_copy":
            prompt = f"请写一篇小红书风格的种草文案。\n产品/主题：{f.get('topic', '')}\n卖点/内容：{f.get('content', '')}\n\n要求：标题吸引人（带emoji），正文亲切活泼，包含Tag。"
        elif request.task == "video_script":
            prompt = f"请写一份短视频脚本。\n主题：{f.get('topic', '')}\n\n要求：包含分镜描述、台词、画面建议，时长约1分钟。"

        # --- Life ---
        elif request.task == "recipe_gen":
            prompt = f"请根据以下食材生成一份菜谱：\n食材：{f.get('content', '')}\n\n要求：包含菜名、所需配料、详细烹饪步骤。"
        elif request.task == "travel_plan":
            prompt = f"请为我制定一份旅行计划。\n目的地：{f.get('destination', '')}\n天数：{f.get('days', '')}\n\n要求：包含每日行程安排、景点推荐、美食建议。"

        # --- Career ---
        elif request.task == "resume_polish":
            prompt = f"请帮我优化这份简历内容：\n{f.get('content', '')}\n\n要求：使用专业职场术语，突出成就和数据，优化排版建议。"
        elif request.task == "interview_prep":
            prompt = f"即将面试岗位：{f.get('topic', '')}\n\n请列出5个高频面试题，并给出优秀的回答思路（STAR法则）。"

        # --- Business ---
        elif request.task == "swot_analysis":
            prompt = f"请对以下项目/主题进行SWOT分析：\n{f.get('content', '')}\n\n要求：列出优势(S)、劣势(W)、机会(O)、威胁(T)，并给出战略建议。"
        elif request.task == "contract_review":
            prompt = f"我是乙方/个人，请帮我审查虽然条款，指出潜在风险和陷阱：\n{f.get('content', '')}\n\n要求：通俗易懂，标出高风险条款。"

        # --- Writing ---
        elif request.task == "title_gen":
            prompt = f"请为这篇文章生成10个爆款标题：\n内容/主题：{f.get('content', '')}\n目标受众：{f.get('topic', '')}\n\n要求：包含数字、悬念、痛点等爆款元素。"
        elif request.task == "article_polish":
            prompt = f"请润色以下文章，使其更流畅、更有文采：\n{f.get('content', '')}\n\n要求：修正语病，提升词汇丰富度，保持原意。"
        
        # --- Excel Generation ---
        elif request.task == "excel_gen":
            prompt = f"请根据以下描述生成一份CSV格式的数据，用于Excel处理：\n描述：{f.get('content', '')}\n\n要求：只输出CSV内容，不要有任何解释，首行为表头。"
            # Handle non-streaming response for Excel
            try:
                def process_excel(csv_content):
                    # Cleanup potential markdown code blocks
                    if "```csv" in csv_content:
                        csv_content = csv_content.split("```csv")[1].split("```")[0].strip()
                    elif "```" in csv_content:
                        csv_content = csv_content.split("```")[1].split("```")[0].strip()
                    
                    # Convert to Excel
                    df = pd.read_csv(io.StringIO(csv_content), on_bad_lines='warn', sep=None, engine='python')
                    
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                        tmp_path = tmp.name
                    
                    df.to_excel(tmp_path, index=False)
                    
                    file_id = str(uuid.uuid4())
                    TEMP_DOWNLOADS[file_id] = {
                        "path": tmp_path,
                        "filename": f"generated_excel_{int(time.time())}.xlsx",
                        "media_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        "ts": time.time()
                    }
                    return {"download_url": f"/api/download/{file_id}"}
                
                return StreamingResponse(
                    keepalive_llm_stream(
                        provider=request.provider,
                        model=resolve_model(request.model),
                        api_key=resolve_api_key(request.api_key),
                        messages=[{"role": "user", "content": prompt}],
                        process_fn=process_excel,
                        **extra
                    ),
                    media_type="application/json"
                )
            except Exception as e:
                print(f"DEBUG EXCEL GEN ERROR: {e}")
                import traceback
                traceback.print_exc()
                raise HTTPException(status_code=500, detail=f"Excel Gen Failed: {str(e)}")

        else:
            prompt = f"请完成以下任务：{f.get('content', '')}"

        async def event_generator():
            try:
                system_instruction = "你是一个专业的AI内容生成助手。请直接根据用户的要求生成最终的内容草稿，不要追问细节，不要反问。如果信息不足，请自行补充合理的假设内容以完成生成。"
                messages = [
                    {"role": "system", "content": system_instruction},
                    {"role": "user", "content": prompt}
                ]
                async for chunk in acall_llm_stream(
                    provider=request.provider,
                    model=resolve_model(request.model),
                    api_key=resolve_api_key(request.api_key),
                    messages=messages,
                    **extra
                ):
                    yield chunk
            except Exception as e:
                yield f"Error: {str(e)}"

        return StreamingResponse(event_generator(), media_type="text/plain")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/generate/creative/upload")
async def generate_creative_upload(
    file: UploadFile = File(...),
    task: str = Form(...),
    content: str = Form(""),
    provider: str = Form(...),
    model: str = Form(...),
    api_key: str = Form(""),
    base_url: str = Form(None),
    receiver: str = Form(None),
    topic: str = Form(None),
    target_lang: str = Form(None),
    time: str = Form(None),
    destination: str = Form(None),
    days: str = Form(None)
):
    try:
        # 1. Parse File
        file_content = ""
        try:
            file_bytes = await file.read()
            # Mock a file-like object with name attribute for parser
            class NamedBytesIO(io.BytesIO):
                def __init__(self, content, name):
                    super().__init__(content)
                    self.name = name
            
            f_obj = NamedBytesIO(file_bytes, file.filename)
            file_content = parse_file(f_obj)
        except Exception as e:
            return JSONResponse(status_code=400, content={"detail": f"File Parse Error: {str(e)}"})

        # 2. Append file content to description
        full_content = f"{content}\n\n【附件内容 ({file.filename})】:\n{file_content}"
        
        # 3. Construct Request Object
        fields = {
            "content": full_content,
            "receiver": receiver,
            "topic": topic,
            "target_lang": target_lang,
            "time": time,
            "destination": destination,
            "days": days
        }
        
        # Filter None values
        fields = {k: v for k, v in fields.items() if v is not None}

        req = CreativeRequest(
            task=task,
            fields=fields,
            provider=provider,
            model=model,
            api_key=resolve_api_key(api_key),
            base_url=base_url
        )
        
        # 4. Delegate to existing logic
        return await generate_creative(req)

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/code/generate")
async def generate_code(request: CodeRequest):
    try:
        role_map = {
            "generate": "你是一个资深程序员。请根据需求生成高质量、可运行的代码。",
            "review": "你是一个代码审查专家。请Review以下代码，指出潜在Bug、性能问题和改进建议。",
            "debug": "你是一个Debug专家。请分析以下代码的错误，并给出修复后的代码。",
            "explain": "你是一个计算机科学教师。请通俗易懂地解释以下代码的逻辑和功能。"
        }
        
        system_role = role_map.get(request.task, "你是一个全能编程助手。")
        user_content = f"语言：{request.language}\n\n内容/需求：\n{request.content}"
        
        messages = [
            {"role": "system", "content": system_role},
            {"role": "user", "content": user_content}
        ]

        extra = {}
        if request.base_url: extra["api_base"] = request.base_url

        async def event_generator():
            try:
                async for chunk in acall_llm_stream(
                    provider=request.provider,
                    model=resolve_model(request.model),
                    api_key=resolve_api_key(request.api_key),
                    messages=messages,
                    **extra
                ):
                    yield chunk
            except Exception as e:
                yield f"Error: {str(e)}"

        return StreamingResponse(event_generator(), media_type="text/plain")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/excel/process")
async def process_excel(
    file: UploadFile = File(...),
    instruction: str = Form(...),
    provider: str = Form("OpenRouter"),
    model: str = Form("gpt-3.5-turbo"),
    api_key: str = Form(""),
    base_url: str = Form(None)
):
    try:
        api_key = resolve_api_key(api_key)
        # 1. Save uploaded file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx", mode="wb") as input_tmp:
            input_path = input_tmp.name
            content = await file.read()
            input_tmp.write(content)

        # 2. Analyze Excel Structure
        try:
            xl = pd.ExcelFile(input_path)
            sheet_names = xl.sheet_names
            schema_info = []
            for sheet in sheet_names[:3]: # Limit to first 3 sheets to avoid token overflow
                df_preview = pd.read_excel(input_path, sheet_name=sheet, nrows=5)
                columns = df_preview.columns.tolist()
                schema_info.append(f"Sheet '{sheet}': Columns {columns}")
            schema_str = "\n".join(schema_info)
        except Exception as e:
            return JSONResponse(status_code=400, content={"detail": f"无法读取 Excel 文件: {str(e)}"})

        # 3. Construct Prompt for Code Generation
        output_path = input_path.replace(".xlsx", "_processed.xlsx")
        
        system_prompt = (
            "你是一个 Python 数据分析专家。你的任务是编写 Python 代码来处理 Excel 文件。\n"
            "【规则】\n"
            "1. 使用 pandas 和 openpyxl 库。\n"
            "2. 代码必须是完整的、可执行的脚本。\n"
            "3. 输入文件路径为变量 `INPUT_PATH`，输出文件路径为变量 `OUTPUT_PATH`。\n"
            "4. 不要使用 input()，不要有交互式命令。\n"
            "5. 只输出代码，不要包含任何 Markdown 标记（如 ```python），直接输出代码内容。\n"
            "6. 必须处理并保存文件到 `OUTPUT_PATH`。\n"
        )
        
        user_prompt = (
            f"输入文件: {input_path}\n"
            f"输出文件: {output_path}\n"
            f"文件结构:\n{schema_str}\n\n"
            f"用户需求: {instruction}\n\n"
            "请编写 Python 代码实现上述需求。代码中直接使用给定的路径变量。"
        )

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt}
        ]

        # 4. Call LLM to Generate Code
        extra = {}
        if base_url: extra["api_base"] = base_url
        
        def process_code_and_run(generated_code):
            # 5. Clean Code
            code = generated_code.strip()
            if "```python" in code:
                code = code.split("```python")[1].split("```")[0]
            elif "```" in code:
                code = code.split("```")[1].split("```")[0]
            
            # Inject path variables to ensure safety/correctness regardless of LLM output
            path_injection = f"INPUT_PATH = r'{input_path}'\nOUTPUT_PATH = r'{output_path}'\n"
            final_code = path_injection + code

            # 6. Execute Code
            print("DEBUG: Executing Excel Processing Code...")
            # Save code to temp file for debugging/execution
            with tempfile.NamedTemporaryFile(delete=False, suffix=".py", mode="w", encoding="utf-8") as code_tmp:
                code_tmp.write(final_code)
                code_path = code_tmp.name
                
            try:
                result = subprocess.run(
                    [sys.executable, code_path],
                    capture_output=True,
                    text=True,
                    timeout=60 # 60s timeout
                )
                
                if result.returncode != 0:
                    print(f"Execution Error: {result.stderr}")
                    raise Exception(f"代码执行失败:\n{result.stderr}\n\n生成的代码:\n{final_code}")
                    
            except subprocess.TimeoutExpired:
                raise Exception("处理超时 (60s)")
            except Exception as e:
                raise e

            # 7. Return Result
            if os.path.exists(output_path):
                file_id = str(uuid.uuid4())
                filename = f"processed_{int(time.time())}.xlsx"
                TEMP_DOWNLOADS[file_id] = {
                    "path": output_path,
                    "filename": filename,
                    "media_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    "ts": time.time()
                }
                return {"download_url": f"/api/download/{file_id}"}
            else:
                raise Exception("输出文件未生成。可能是代码逻辑错误。")

        return StreamingResponse(
            keepalive_llm_stream(
                provider=provider,
                model=model,
                api_key=resolve_api_key(api_key),
                messages=messages,
                process_fn=process_code_and_run,
                **extra
            ),
            media_type="application/json"
        )

    except Exception as e:
        traceback.print_exc()
        return JSONResponse(status_code=500, content={"detail": str(e)})


@app.post("/api/convert/img-to-excel")
async def img_to_excel(file: UploadFile = File(...)):
    try:
        from PIL import Image
        import pandas as pd
        # In a real scenario, we would use a Vision LLM (like Gemini Pro Vision or GPT-4o) 
        # to extract data from the image. 
        # Since we don't have a direct "Image to Excel" library that works perfectly without OCR/AI,
        # we will simulate a success or use a placeholder if we can't call an AI here easily for file uploads 
        # without user's API key in this context. 
        # However, for this task, the user wants "various file to excel". 
        # We will assume a basic implementation or placeholder for now, 
        # OR we can try to use standard OCR if available, but requirements.txt didn't include tesseract.
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/persona/analyze")
async def persona_analyze(
    file: UploadFile = File(...),
    provider: str = Form("OpenRouter"),
    model: str = Form("gpt-3.5-turbo"),
    api_key: str = Form(""),
    base_url: str = Form(None)
):
    try:
        api_key = resolve_api_key(api_key)
        # 1. Parse File
        file_bytes = await file.read()
        
        class NamedBytesIO(io.BytesIO):
            def __init__(self, content, name):
                super().__init__(content)
                self.name = name
        
        f_obj = NamedBytesIO(file_bytes, file.filename)
        file_content = parse_file(f_obj)
        
        if file_content.startswith("["): # Check for parser errors
             return JSONResponse(status_code=400, content={"detail": file_content})

        # 2. Analyze Style
        system_prompt = await analyze_chat_style(
            text_content=file_content,
            provider=provider,
            model=model,
            api_key=api_key,
            base_url=base_url
        )
        
        return {"system_prompt": system_prompt}

    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/api/persona/local_scan")
async def persona_local_scan():
    try:
        files = scan_qq_logs()
        return {"files": files}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/persona/analyze_text")
async def persona_analyze_text(
    text: str = Form(...),
    provider: str = Form("OpenRouter"),
    model: str = Form("gpt-3.5-turbo"),
    api_key: str = Form(""),
    base_url: str = Form(None)
):
    try:
        api_key = resolve_api_key(api_key)
        system_prompt = await analyze_chat_style(
            text_content=text,
            provider=provider,
            model=model,
            api_key=resolve_api_key(api_key),
            base_url=base_url
        )
        return {"system_prompt": system_prompt}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/persona/import_local")
async def persona_import_local(
    path: str = Form(...),
    provider: str = Form("OpenRouter"),
    model: str = Form("gpt-3.5-turbo"),
    api_key: str = Form(""),
    base_url: str = Form(None)
):
    try:
        api_key = resolve_api_key(api_key)
        if not os.path.exists(path):
            raise HTTPException(status_code=404, detail="File not found")
            
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            file_content = f.read()
            
        system_prompt = await analyze_chat_style(
            text_content=file_content,
            provider=provider,
            model=model,
            api_key=api_key,
            base_url=base_url
        )
        return {"system_prompt": system_prompt}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))
        # Let's implement a mockup that returns a sample Excel for demonstration 
        # OR if the user provides a CSV/Excel-like image, we can't easily parse it without AI.
        
        # ACTUALLY, the user's request "various file convert to excel" might imply extracting tables.
        # Let's create a dummy Excel file for now to demonstrate the FLOW, 
        # as integrating Tesseract/Vision LLM requires more setup (API Key handling for file uploads).
        # Wait, we can't easily pass the API key in a file upload form without changing the frontend to send it.
        # For this step, I will implement a "PDF/Image to Excel" using basic libraries if possible, or Mock it.
        # Given the "AI Mate" context, it likely expects AI. 
        # But for "Format Conversion", usually it's structural.
        # Let's stick to valid conversions. 
        # PDF to Excel is possible with `pdfplumber`.
        # Image to Excel is hard without OCR.
        
        # Let's implement PDF to Excel using pdfplumber.
        return JSONResponse(status_code=400, content={"detail": "Image to Excel requires OCR/AI integration. Please use PDF to Excel for now."})

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/convert/pdf-to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    try:
        import pdfplumber
        import pandas as pd
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
             tmp_pdf.write(await file.read())
             tmp_pdf_path = tmp_pdf.name
        
        excel_path = tmp_pdf_path.replace(".pdf", ".xlsx")
        
        try:
            with pdfplumber.open(tmp_pdf_path) as pdf:
                all_tables = []
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        all_tables.append(df)
            
            if not all_tables:
                 raise HTTPException(status_code=400, detail="No tables found in PDF")
                 
            # Save to Excel
            with pd.ExcelWriter(excel_path) as writer:
                for i, df in enumerate(all_tables):
                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
                    
            return FileResponse(
                excel_path, 
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", 
                filename=f"{file.filename.replace('.pdf', '')}.xlsx"
            )
        finally:
            if os.path.exists(tmp_pdf_path):
                os.remove(tmp_pdf_path)
            # We don't remove excel_path immediately so it can be streamed, 
            # but usually we should clean up. Windows file locking might be an issue.
            # For this simple server, we rely on OS to clean temp eventually or background task.
            
    except HTTPException:
        raise
    except ImportError:
        raise HTTPException(status_code=500, detail="pdfplumber/pandas not installed")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/api/convert/pdf-to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    try:
        from pdf2docx import Converter
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
             tmp_pdf.write(await file.read())
             tmp_pdf_path = tmp_pdf.name
        
        word_path = tmp_pdf_path.replace(".pdf", ".docx")
        
        try:
            cv = Converter(tmp_pdf_path)
            cv.convert(word_path, start=0, end=None)
            cv.close()
            
            return FileResponse(
                word_path, 
                media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document", 
                filename=f"{file.filename.replace('.pdf', '')}.docx"
            )
        finally:
            if os.path.exists(tmp_pdf_path):
                os.remove(tmp_pdf_path)
            
    except ImportError:
        raise HTTPException(status_code=500, detail="pdf2docx not installed")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/convert/img-to-pdf")
async def img_to_pdf(files: list[UploadFile] = File(...)):
    try:
        from PIL import Image
        
        if not files:
             raise HTTPException(status_code=400, detail="No files uploaded")
             
        images = []
        for file in files:
            content = await file.read()
            img = Image.open(io.BytesIO(content))
            if img.mode == "RGBA":
                img = img.convert("RGB")
            images.append(img)
            
        if not images:
             raise HTTPException(status_code=400, detail="No valid images found")
             
        pdf_bytes = io.BytesIO()
        images[0].save(
            pdf_bytes, "PDF", resolution=100.0, save_all=True, append_images=images[1:]
        )
        pdf_bytes.seek(0)
        
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_pdf:
            tmp_pdf.write(pdf_bytes.getvalue())
            tmp_pdf_path = tmp_pdf.name
            
        return FileResponse(tmp_pdf_path, media_type="application/pdf", filename="merged_images.pdf")

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate/ppt/outline")
async def generate_ppt_outline(request: PPTRequest):
    try:
        _PPT_PROMPT = """
        你是一个专业的 PPT 设计师。
        请根据用户的主题："{topic}"
        生成一份 PPT 大纲，包含 5-8 页内容。
        
        请严格返回合法的 JSON 格式，结构如下：
        {{
          "slides": [
            {{
              "title": "页标题",
              "content": ["要点1", "要点2", "要点3"]
            }},
            ...
          ]
        }}
        
        注意：
        1. 第一页必须是封面页（"title": "主题", "content": ["副标题/演讲者"]）。
        2. 内容要精炼、专业。
        3. 不要包含 JSON 以外的任何文字。
        """
        
        prompt = _PPT_PROMPT.format(topic=request.topic)
        messages = [{"role": "user", "content": prompt}]
        
        extra = {}
        if request.base_url:
            extra["api_base"] = request.base_url

        def process_ppt(response_text):
            clean_json = re.sub(r"```json\s*", "", response_text)
            clean_json = re.sub(r"\s*```", "", clean_json)
            data = json.loads(clean_json)
            return {"data": data}

        return StreamingResponse(
            keepalive_llm_stream(
                provider=request.provider,
                model=resolve_model(request.model),
                api_key=resolve_api_key(request.api_key),
                messages=messages,
                process_fn=process_ppt,
                **extra
            ),
            media_type="application/json"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/generate/ppt/create")
async def create_ppt(request: PPTCreateRequest):
    try:
        from pptx import Presentation
    except ImportError:
        raise HTTPException(status_code=500, detail="python-pptx not installed")

    data = request.data
    slides_data = data.get("slides", [])

    prs = Presentation()

    if slides_data:
        layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        cover = slides_data[0]
        title.text = cover.get("title", "Presentation")
        subtitle.text = "\n".join(cover.get("content", []))

    layout = prs.slide_layouts[1]
    for slide_info in slides_data[1:]:
        slide = prs.slides.add_slide(layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        title_shape.text = slide_info.get("title", "Untitled")
        
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        points = slide_info.get("content", [])
        if points:
            tf.text = points[0]
            for point in points[1:]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 1
                
    ppt_buffer = io.BytesIO()
    prs.save(ppt_buffer)
    ppt_buffer.seek(0)
    
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_ppt:
        tmp_ppt.write(ppt_buffer.getvalue())
        tmp_path = tmp_ppt.name
        
    return FileResponse(
        tmp_path, 
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="generated_presentation.pptx"
    )

@app.post("/api/analyze/table")
async def analyze_table(file: UploadFile = File(...)):
    try:
        import pandas as pd
        content = await file.read()
        try:
            if file.filename.endswith(".csv"):
                df = pd.read_csv(io.BytesIO(content))
            else:
                df = pd.read_excel(io.BytesIO(content))
        except Exception:
             raise HTTPException(status_code=400, detail="Failed to parse file. Please upload a valid CSV or Excel file.")
        
        df = df.where(pd.notnull(df), None)
        columns = df.columns.tolist()
        numeric_cols = df.select_dtypes(include=['number']).columns.tolist()
        preview_data = df.head(100).to_dict(orient="records")
        
        return {
            "filename": file.filename,
            "columns": columns,
            "numeric_columns": numeric_cols,
            "data": preview_data,
            "total_rows": len(df)
        }

    except ImportError:
        raise HTTPException(status_code=500, detail="pandas or openpyxl not installed")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class MindMapRequest(BaseModel):
    topic: str
    chart_type: str = "mindmap" # mindmap, flowchart, timeline, gantt, sequence, class, state, pie
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None

@app.post("/api/generate/mindmap")
async def generate_mindmap(request: MindMapRequest):
    try:
        # Prompts for different chart types
        prompts = {
            "mindmap": """
                你是一个思维导图专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的思维导图代码 (mindmap)。
                
                【输出规则】
                1. 必须以 `mindmap` 关键字开头。
                2. 根节点格式：`root((主题名))`。
                3. 使用空格缩进表示层级，不要使用 Tab。
                4. 仅返回代码本身，禁止包含任何 Markdown 代码块标记（如 ```mermaid）或解释文字。
                5. 至少扩展 3 层，包含 10-15 个节点。
                6. 节点文字如果包含括号或特殊字符，请用双引号包裹，例如：`"分支(备注)"`。
                """,
            "flowchart": """
                你是一个流程图专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的流程图代码 (graph TD)。
                
                【输出规则】
                1. 必须包含 **分支判断逻辑 (Decision Nodes)**，使用 `{{}}` 语法（如：`A{{是否成功}} -- 是 --> B`）。
                2. 流程必须是 **非线性的**，包含多个分支、交汇点或循环（Feedback Loop）。
                3. 节点文字如果包含特殊字符，请务必用双引号包裹，例如：`id1["文字内容"]`。
                4. 包含至少 10-15 个具有逻辑意义的节点。
                5. 仅返回 Mermaid 代码，禁止包含 Markdown 代码块标记或解释文字。
                6. 必须以 `graph TD` 开头。
                """,
            "timeline": """
                你是一个时间序列专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的时间轴代码 (timeline)。
                
                要求：
                1. 包含至少 5-8 个关键时间点。
                2. 仅返回 Mermaid 代码，禁止包含 Markdown 代码块标记或解释文字。
                3. 必须以 `timeline` 开头。
                4. 格式：`年 : 事件` 或 `时间段 : 事件`。
                """,
            "gantt": """
                你是一个项目规划专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的甘特图代码 (gantt)。
                
                要求：
                1. 包含日期格式 `dateFormat YYYY-MM-DD`。
                2. 包含 `title` 关键字定义标题。
                3. 仅返回 Mermaid 代码，禁止包含 Markdown 代码块标记或解释文字。
                4. 必须以 `gantt` 开头。
                """,
            "sequence": """
                你是一个时序图专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的时序图代码 (sequenceDiagram)。
                
                要求：
                1. 描述参与者之间的交互，包含至少 3 个参与者。
                2. 仅返回 Mermaid 代码，禁止包含 Markdown 代码块标记或解释文字。
                3. 必须以 `sequenceDiagram` 开头。
                """,
             "class": """
                你是一个架构专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的类图代码 (classDiagram)。
                
                要求：
                1. 定义类及其属性、方法和关系。
                2. 仅返回 Mermaid 代码，禁止包含 Markdown 代码块标记或解释文字。
                3. 必须以 `classDiagram` 开头。
                """,
             "state": """
                你是一个状态机专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的状态图代码 (stateDiagram-v2)。
                
                要求：
                1. 描述状态转换流程。
                2. 仅返回 Mermaid 代码，禁止包含 Markdown 代码块标记或解释文字。
                3. 必须以 `stateDiagram-v2` 开头。
                """,
             "pie": """
                你是一个数据专家。请根据用户的主题："{topic}"
                生成一份 **Mermaid JS** 的标准饼图代码 (pie)。
                
                【严等规则】
                1. 必须以 `pie` 开头。
                2. 可以包含标题，格式为 `title 标题文字`。
                3. 数据行格式：`"标签" : 数值`。**标签必须用双引号包裹**。
                4. 禁止包含任何非 Mermaid 语法的解释说明。
                5. 仅返回 Mermaid 代码。

                【示例】
                pie title 饼图标题
                    "分类A" : 40
                    "分类B" : 60
                """
        }
        
        selected_prompt = prompts.get(request.chart_type, prompts["mindmap"])
        prompt = selected_prompt.format(topic=request.topic)
        
        messages = [{"role": "user", "content": prompt}]
        
        extra = {}
        if request.base_url:
            extra["api_base"] = request.base_url

        def process_mindmap(response_text):
            code = response_text.strip()
            # 1. Handle Markdown Code Blocks
            blocks = re.findall(r"```(?:mermaid)?\s*(.*?)\s*```", code, re.DOTALL)
            if blocks:
                code = blocks[0].strip()
            # 2. Extract only the mermaid part (remove AI chatter)
            start_keywords = ["mindmap", "graph", "flowchart", "timeline", "gantt", "sequenceDiagram", "classDiagram", "stateDiagram", "pie"]
            lines = code.split('\n')
            valid_start = -1
            for i, line in enumerate(lines):
                clean_line = line.strip()
                if any(clean_line.startswith(kw) for kw in start_keywords):
                    valid_start = i
                    break
            if valid_start != -1:
                code = "\n".join(lines[valid_start:]).strip()
            # 3. Specific Fixes for Pie Charts
            if request.chart_type == "pie":
                if not code.startswith("pie"):
                    code = "pie\n" + code
                code = re.sub(r'^\s*([^"\s:][^:\n]+[^"\s:])\s*:\s*(\d+)', r'    "\1" : \2', code, flags=re.MULTILINE)
            # 4. Specific Fixes for Mindmaps
            if request.chart_type == "mindmap" and not code.startswith("mindmap"):
                code = "mindmap\n  " + code.replace("\n", "\n  ")
            return {"code": code}

        return StreamingResponse(
            keepalive_llm_stream(
                provider=request.provider,
                model=resolve_model(request.model),
                api_key=resolve_api_key(request.api_key),
                messages=messages,
                process_fn=process_mindmap,
                **extra
            ),
            media_type="application/json"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/system/generate_code")
async def system_generate_code(request: SystemRequest):
    try:
        import platform
        _SYSTEM_PROMPT = """
        你是一个精通 Python 系统操作的 AI 助手。
        用户的需求是：{query}
        
        请生成一段 Python 代码来实现该需求。
        要求：
        1. 仅使用 Python 标准库（如 os, shutil, glob, pathlib, datetime 等），严禁使用第三方库。
        2. 代码必须是完整的、可执行的脚本。
        3. 代码中不要包含 input() 等交互函数。
        4. 将代码包裹在 ```python ... ``` 代码块中。
        5. 不要解释代码，只返回代码块。
        
        当前工作目录：{cwd}
        操作系统：{os_name}
        """
        
        prompt = _SYSTEM_PROMPT.format(
            query=request.query,
            cwd=os.getcwd(),
            os_name=platform.system()
        )
        
        messages = [{"role": "user", "content": prompt}]
        
        extra = {}
        if request.base_url:
            extra["api_base"] = request.base_url

        def process_code(response_text):
            code = re.sub(r"```python(.*?)```", r"\1", response_text, flags=re.DOTALL)
            code = re.sub(r"```", "", code).strip()
            return {"code": code}

        return StreamingResponse(
            keepalive_llm_stream(
                provider=request.provider,
                model=resolve_model(request.model),
                api_key=resolve_api_key(request.api_key),
                messages=messages,
                process_fn=process_code,
                **extra
            ),
            media_type="application/json"
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/api/system/execute_code")
async def system_execute_code(request: SystemExecuteRequest):
    try:
        import sys
        import io
        from contextlib import redirect_stdout
        import traceback
        
        code = request.code
        buffer = io.StringIO()
        
        try:
             with redirect_stdout(buffer):
                exec(code, {"__name__": "__main__"})
             output = buffer.getvalue()
        except Exception:
             output = buffer.getvalue() + "\n" + traceback.format_exc()
             
        return {"output": output}

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

class GameAdventureRequest(BaseModel):
    world_setting: str
    current_state: dict
    history: list[dict]
    user_action: str
    provider: str = "OpenRouter"
    model: str = "gpt-3.5-turbo"
    api_key: str = ""
    base_url: str | None = None


@app.post("/api/game/adventure")
async def game_adventure(request: GameAdventureRequest):
    try:
        # Construct System Prompt for the AI DM
        setting = request.world_setting
        state = request.current_state
        action = request.user_action
        
        system_prompt = f"""You are the Dungeon Master (DM) for a text-based RPG.
World Setting: {setting}

Player State:
- HP: {state.get('hp', 100)}/{state.get('max_hp', 100)}
- Inventory: {', '.join(state.get('inventory', [])) or 'Empty'}
- Location: {state.get('location', 'Unknown')}
- Status: {state.get('status', 'Normal')}

Current Situation:
The player performs the action: "{action}"

Your Task:
1. Narrate the result of the action and the new situation. 
   - **CRITICAL:** Ensure the narrative is logically consistent with the previous state. Avoid sudden jumps in time or location unless the action explicitly requires it. 
   - Focus on cause and effect. 
   - Be descriptive and immersive, but keep the story grounded in the current context.
2. Update the player's state based on the events (e.g., if they drink a potion, HP increases; if they find an item, add to inventory).
3. Offer 3 distinct choices for what the player might do next (to help them, but they are free to do anything).

Output JSON format ONLY:
{{
    "plot": "The narrative text...",
    "state_update": {{
        "hp": 100,             // Current HP
        "inventory": [],       // Updated list
        "location": "...",     // Current location name
        "status": "..."        // Short status description
    }},
    "choices": ["Choice 1", "Choice 2", "Choice 3"]
}}
"""
        # Limit history to last 6 turns to keep context window manageable but relevant
        recent_history = request.history[-6:] if request.history else []
        
        messages = [{"role": "system", "content": system_prompt}]
        messages.extend(recent_history)
        messages.append({"role": "user", "content": action})

        extra = {}
        if request.base_url:
            extra["api_base"] = request.base_url

        def process_adventure(response_content):
            try:
                json_match = re.search(r"\{.*\}", response_content, re.DOTALL)
                if json_match:
                    result = json.loads(json_match.group(0))
                else:
                    result = json.loads(response_content)
                
                # Robustness: Check for missing keys
                if "plot" not in result:
                    result["plot"] = result.get("narrative") or result.get("story") or result.get("description") or result.get("output") or response_content
                
                if "state_update" not in result:
                    result["state_update"] = {}

                return result
            except json.JSONDecodeError:
                return {
                    "plot": response_content + "\n(系统提示: AI 返回格式异常，但剧情继续)",
                    "state_update": state,
                    "choices": ["继续", "观察四周", "检查状态"]
                }

        return StreamingResponse(
            keepalive_llm_stream(
                provider=request.provider,
                model=resolve_model(request.model),
                api_key=resolve_api_key(request.api_key),
                messages=messages,
                process_fn=process_adventure,
                **extra
            ),
            media_type="application/json"
        )

    except Exception as e:
        print(f"Adventure Error: {e}")
        return JSONResponse(status_code=500, content={"detail": str(e)})

# Mount Static Files (Frontend)
app.mount("/", StaticFiles(directory="static", html=True), name="static")

if __name__ == "__main__":
    uvicorn.run("api_server:app", host="0.0.0.0", port=8000, reload=True)
